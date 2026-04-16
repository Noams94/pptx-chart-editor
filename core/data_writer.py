"""Write edited DataFrames back into PPTX charts."""

from __future__ import annotations

from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.oxml.ns import qn
from lxml import etree

from openpyxl import load_workbook

from core.data_extractor import is_percentage_format


def _display_to_raw(df: pd.DataFrame, series_formats: dict) -> pd.DataFrame:
    """Convert display values back to raw values.

    E.g., 67 (displayed as 67%) -> 0.67 (raw value stored in chart)
    """
    raw_df = df.copy()
    for col in df.columns[1:]:  # Skip first column (categories)
        fmt = series_formats.get(col, "General")
        if is_percentage_format(fmt):
            raw_df[col] = df[col].apply(
                lambda v: v / 100.0 if pd.notna(v) else None
            )
    return raw_df


def update_chart_data(
    pptx_bytes: bytes,
    slide_index: int,
    shape_name: str,
    display_df: pd.DataFrame,
    is_xy: bool = False,
    series_formats: dict = None,
    series_visibility: dict = None,
    shape_id: int = None,
    series_colors: dict = None,
) -> bytes:
    """Update a single chart's data in the PPTX and return updated bytes.

    Args:
        pptx_bytes: Original PPTX file bytes
        slide_index: Zero-based slide index
        shape_name: Name of the chart shape
        display_df: DataFrame with display values (67 for 67%)
        is_xy: Whether this is an XY/scatter chart
        series_formats: Dict of series_name -> formatCode for converting display->raw

    Returns:
        Updated PPTX file bytes
    """
    # Convert display values to raw values
    if series_formats:
        df = _display_to_raw(display_df, series_formats)
    else:
        df = display_df

    prs = Presentation(BytesIO(pptx_bytes))
    slide = prs.slides[slide_index]

    # Find the chart shape
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart and shape.name == shape_name:
            chart_shape = shape
            break

    if chart_shape is None:
        raise ValueError(f"Chart '{shape_name}' not found on slide {slide_index + 1}")

    chart = chart_shape.chart

    if is_xy:
        chart_data = XyChartData()
        col_names = df.columns.tolist()
        for i in range(0, len(col_names), 2):
            series_name = col_names[i].replace("X_", "")
            series = chart_data.add_series(series_name)
            x_vals = df.iloc[:, i].dropna().tolist()
            y_vals = df.iloc[:, i + 1].dropna().tolist()
            for x, y in zip(x_vals, y_vals):
                series.add_data_point(x, y)
    else:
        chart_data = CategoryChartData()
        categories = df.iloc[:, 0].dropna().astype(str).tolist()
        chart_data.categories = categories

        for col in df.columns[1:]:
            values = df[col].tolist()
            values = [None if pd.isna(v) else float(v) for v in values]
            values = values[:len(categories)]
            while len(values) < len(categories):
                values.append(None)
            chart_data.add_series(col, values)

    chart.replace_data(chart_data)

    # Restore number format codes that replace_data() resets to "General"
    if series_formats:
        _restore_format_codes(chart, series_formats)
        _format_embedded_excel(chart, series_formats)

    # Restore series visibility state
    if series_visibility:
        _restore_visibility(chart, series_visibility)

    # Restore series colors
    if series_colors:
        _restore_series_colors(chart, series_colors)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def update_multiple_charts(
    pptx_bytes: bytes,
    updates: list[tuple[int, str, pd.DataFrame, bool, dict | None, dict | None]],
) -> bytes:
    """Update multiple charts in a single parse/save cycle.

    Each update is a tuple: (slide_index, shape_name, display_df, is_xy, series_formats, series_visibility)
    """
    prs = Presentation(BytesIO(pptx_bytes))

    for update in updates:
        # Support 5-tuple, 6-tuple, and 7-tuple (with shape_id) formats
        if len(update) == 7:
            slide_index, shape_name, display_df, is_xy, series_formats, series_visibility, _shape_id = update
        elif len(update) == 6:
            slide_index, shape_name, display_df, is_xy, series_formats, series_visibility = update
        else:
            slide_index, shape_name, display_df, is_xy, series_formats = update
            series_visibility = None
        df = _display_to_raw(display_df, series_formats) if series_formats else display_df
        slide = prs.slides[slide_index]

        chart_shape = None
        for shape in slide.shapes:
            if shape.has_chart and shape.name == shape_name:
                chart_shape = shape
                break
        if chart_shape is None:
            continue

        chart = chart_shape.chart

        if is_xy:
            chart_data = XyChartData()
            col_names = df.columns.tolist()
            for i in range(0, len(col_names), 2):
                series_name = col_names[i].replace("X_", "")
                series = chart_data.add_series(series_name)
                x_vals = df.iloc[:, i].dropna().tolist()
                y_vals = df.iloc[:, i + 1].dropna().tolist()
                for x, y in zip(x_vals, y_vals):
                    series.add_data_point(x, y)
        else:
            chart_data = CategoryChartData()
            categories = df.iloc[:, 0].dropna().astype(str).tolist()
            chart_data.categories = categories
            for col in df.columns[1:]:
                values = df[col].tolist()
                values = [None if pd.isna(v) else float(v) for v in values]
                values = values[:len(categories)]
                while len(values) < len(categories):
                    values.append(None)
                chart_data.add_series(col, values)

        chart.replace_data(chart_data)
        if series_formats:
            _restore_format_codes(chart, series_formats)
            _format_embedded_excel(chart, series_formats)
        if series_visibility:
            _restore_visibility(chart, series_visibility)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _format_embedded_excel(chart, series_formats: dict):
    """Apply number formats to the embedded Excel workbook inside the chart.

    chart.replace_data() regenerates the embedded xlsx with 'General' format
    on all cells. When PowerPoint refreshes data from this source, it displays
    raw decimals (0.15) instead of formatted percentages (15%).

    This writes the correct number format codes into the embedded workbook
    cells so PowerPoint inherits the intended display format.
    """
    try:
        xlsx_part = chart.part.chart_workbook.xlsx_part
        wb = load_workbook(BytesIO(xlsx_part.blob))
    except Exception:
        return  # No embedded workbook to format

    ws = wb.active
    if ws is None or ws.max_row is None or ws.max_column is None:
        return

    format_values = list(series_formats.values())

    # Column 1 = categories, columns 2+ = series data (1-based indexing)
    for col_offset, fmt_code in enumerate(format_values):
        excel_col = col_offset + 2  # skip category column
        if excel_col > ws.max_column:
            break
        for row in range(2, ws.max_row + 1):  # skip header row
            cell = ws.cell(row=row, column=excel_col)
            if cell.value is not None:
                cell.number_format = fmt_code

    buf = BytesIO()
    wb.save(buf)
    xlsx_part.blob = buf.getvalue()


def _restore_format_codes(chart, series_formats: dict):
    """Restore formatCode in chart XML after replace_data() resets them.

    series_formats maps column_name -> formatCode. We match by index since
    column order matches XML series order.
    """
    chart_xml = chart.part._element

    # Get column names in order (these are the keys of series_formats)
    format_values = list(series_formats.values())

    for idx, ser in enumerate(chart_xml.iter(qn('c:ser'))):
        if idx >= len(format_values):
            break

        fmt_code = format_values[idx]

        # Set formatCode in val > numRef > numCache > formatCode
        val = ser.find(qn('c:val'))
        if val is not None:
            num_ref = val.find(qn('c:numRef'))
            if num_ref is not None:
                num_cache = num_ref.find(qn('c:numCache'))
                if num_cache is not None:
                    fc = num_cache.find(qn('c:formatCode'))
                    if fc is None:
                        fc = etree.SubElement(num_cache, qn('c:formatCode'))
                    fc.text = fmt_code


def _restore_visibility(chart, series_visibility: dict):
    """Restore series visibility (show/hide) in chart XML.

    In PowerPoint, hidden series have <c:delete val="1"/> inside <c:ser>.
    series_visibility maps series_name -> bool (True = visible).
    We match by index since column order matches XML series order.
    """
    chart_xml = chart.part._element
    visibility_values = list(series_visibility.values())

    for idx, ser in enumerate(chart_xml.iter(qn('c:ser'))):
        if idx >= len(visibility_values):
            break

        visible = visibility_values[idx]
        delete_el = ser.find(qn('c:delete'))

        if not visible:
            # Series should be hidden — add or update c:delete val="1"
            if delete_el is None:
                delete_el = etree.SubElement(ser, qn('c:delete'))
            delete_el.set('val', '1')
        else:
            # Series should be visible — remove c:delete if present
            if delete_el is not None:
                ser.remove(delete_el)


def _restore_series_colors(chart, series_colors: dict):
    """Restore series fill colors in chart XML after replace_data() resets them.

    series_colors maps series_name -> hex color string (e.g. '#4472C4') or ''.
    Empty strings are skipped (no color override).
    Sets c:spPr > a:solidFill > a:srgbClr for bar/column/area/pie charts.
    Also updates a:ln > a:solidFill for line charts if a:ln already exists.
    """
    chart_xml = chart.part._element
    color_values = list(series_colors.values())

    for idx, ser in enumerate(chart_xml.iter(qn('c:ser'))):
        if idx >= len(color_values):
            break

        hex_color = color_values[idx]
        if not hex_color:
            continue  # No explicit color — leave as default

        hex_val = hex_color.lstrip('#').upper()

        # Find or create c:spPr
        spPr = ser.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ser, qn('c:spPr'))

        # Set fill solidFill (bar, column, area, pie)
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        else:
            for child in list(solidFill):
                solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', hex_val)

        # Also update line color if a:ln already exists (line/scatter charts)
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            ln_fill = ln.find(qn('a:solidFill'))
            if ln_fill is None:
                ln_fill = etree.SubElement(ln, qn('a:solidFill'))
            else:
                for child in list(ln_fill):
                    ln_fill.remove(child)
            ln_srgb = etree.SubElement(ln_fill, qn('a:srgbClr'))
            ln_srgb.set('val', hex_val)
