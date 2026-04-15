"""Extract chart data from PPTX files into pandas DataFrames."""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ui.rtl_support import t, chart_type_display_name


def _iter_chart_shapes(shapes):
    """Yield all chart-bearing shapes, recursing into group shapes."""
    for shape in shapes:
        if shape.has_chart:
            yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_chart_shapes(shape.shapes)


# Chart types that use XyChartData (scatter plots)
XY_CHART_TYPES = {
    XL_CHART_TYPE.XY_SCATTER,
    XL_CHART_TYPE.XY_SCATTER_LINES,
    XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
}


def is_percentage_format(fmt: str) -> bool:
    """Check if a format code represents percentages (e.g., '0%', '0.0%', '#,##0%')."""
    if not fmt:
        return False
    # Remove escaped characters and quoted strings
    cleaned = re.sub(r'"[^"]*"', '', fmt)
    cleaned = re.sub(r'\\\.', '', cleaned)
    return '%' in cleaned


def _extract_series_formats_by_index(chart: Chart) -> list[str]:
    """Extract number format codes per series from chart XML, ordered by index.

    Returns list of formatCode strings (e.g., ['0%', 'General', '#,##0']),
    one per series in chart order.
    """
    formats = []
    chart_xml = chart.part._element

    # Find the plot group (e.g., c:barChart, c:lineChart) which contains c:ser elements
    # We need to iterate in document order to match series index
    for ser in chart_xml.iter(qn('c:ser')):
        fmt_code = "General"

        # Try val > numRef > numCache > formatCode
        val = ser.find(qn('c:val'))
        if val is not None:
            num_ref = val.find(qn('c:numRef'))
            if num_ref is not None:
                num_cache = num_ref.find(qn('c:numCache'))
                if num_cache is not None:
                    fc = num_cache.find(qn('c:formatCode'))
                    if fc is not None and fc.text:
                        fmt_code = fc.text

        formats.append(fmt_code)

    return formats


def _extract_series_visibility(chart: Chart) -> dict[str, bool]:
    """Extract visibility state per series from chart XML.

    In PowerPoint, hidden series have <c:delete val="1"/> inside <c:ser>.
    Returns dict mapping series index -> visible (True/False).
    """
    visibility = {}
    chart_xml = chart.part._element

    for idx, ser in enumerate(chart_xml.iter(qn('c:ser'))):
        delete_el = ser.find(qn('c:delete'))
        visible = delete_el is None or delete_el.get('val', '0') == '0'
        visibility[idx] = visible

    return visibility


@dataclass
class ChartInfo:
    slide_index: int
    shape_name: str
    shape_id: int  # unique within a slide, needed when shape_name is duplicated
    chart_type: int
    chart_type_name: str
    dataframe: pd.DataFrame          # Display values (67 for 67%)
    is_xy: bool = False
    series_names: list = field(default_factory=list)
    series_formats: dict = field(default_factory=dict)  # series_name -> formatCode
    series_visibility: dict = field(default_factory=dict)  # series_name -> bool (visible)
    chart_title: str = ""            # Chart title from XML (if available)

    @property
    def key(self) -> tuple:
        """Unique identifier for this chart: (slide_index, shape_name, shape_id)."""
        return (self.slide_index, self.shape_name, self.shape_id)


def _extract_chart_data(chart: Chart) -> tuple[pd.DataFrame, bool, list[str], dict, dict]:
    """Extract data from a chart into a display DataFrame."""
    chart_type = chart.chart_type
    is_xy = chart_type in XY_CHART_TYPES

    plot = chart.plots[0]
    series_list = list(plot.series)
    series_names = [s.name if s.name else t("series_n", n=i+1) for i, s in enumerate(series_list)]

    # Extract number formats by index, then map to our column names
    format_list = _extract_series_formats_by_index(chart)
    series_formats = {}
    for i, name in enumerate(series_names):
        if i < len(format_list):
            series_formats[name] = format_list[i]

    # Extract visibility state per series
    visibility_by_idx = _extract_series_visibility(chart)
    series_visibility = {}
    for i, name in enumerate(series_names):
        series_visibility[name] = visibility_by_idx.get(i, True)

    if is_xy:
        data = {}
        for i, series in enumerate(series_list):
            x_vals = list(series.values)
            y_vals = list(series.values)
            data[f"X_{series_names[i]}"] = x_vals
            data[f"Y_{series_names[i]}"] = y_vals
        display_df = pd.DataFrame(data)
    else:
        try:
            categories = [str(c) for c in plot.categories]
        except Exception:
            categories = [str(i + 1) for i in range(len(list(series_list[0].values)))]

        display_data = {t("category"): categories}

        for i, series in enumerate(series_list):
            values = list(series.values)
            while len(values) < len(categories):
                values.append(None)
            values = values[:len(categories)]

            name = series_names[i]
            fmt = series_formats.get(name, "General")
            if is_percentage_format(fmt):
                display_data[name] = [
                    round(v * 100, 2) if v is not None else None
                    for v in values
                ]
            else:
                display_data[name] = values

        display_df = pd.DataFrame(display_data)

    return display_df, is_xy, series_names, series_formats, series_visibility


def extract_all_charts(pptx_bytes: bytes) -> list[ChartInfo]:
    """Extract all charts from a PPTX file."""
    prs = Presentation(BytesIO(pptx_bytes))
    charts = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in _iter_chart_shapes(slide.shapes):
            chart = shape.chart
            try:
                display_df, is_xy, series_names, series_formats, series_visibility = _extract_chart_data(chart)
                # Extract chart title if available
                title_text = ""
                if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
                    title_text = chart.chart_title.text_frame.text.strip()
                info = ChartInfo(
                    slide_index=slide_idx,
                    shape_name=shape.name,
                    shape_id=shape.shape_id,
                    chart_type=chart.chart_type,
                    chart_type_name=chart_type_display_name(chart.chart_type),
                    dataframe=display_df,
                    is_xy=is_xy,
                    series_names=series_names,
                    series_formats=series_formats,
                    series_visibility=series_visibility,
                    chart_title=title_text,
                )
                charts.append(info)
            except Exception as e:
                print(f"Warning: Could not extract chart '{shape.name}' on slide {slide_idx + 1}: {e}")

    return charts
